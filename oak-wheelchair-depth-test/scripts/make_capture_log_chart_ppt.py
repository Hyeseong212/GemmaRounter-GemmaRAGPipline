#!/usr/bin/env python3

from __future__ import annotations

import json
from pathlib import Path


ROOT = Path("/home/rbiotech-server/LLM_Harnes_Support/GemmaRounter-GemmaRAGPipline")
RUNTIME_DATA = ROOT / "oak-wheelchair-depth-test" / "runtime-data"
DEFAULT_LOG = RUNTIME_DATA / "capture_log.json"
DEFAULT_OUT = RUNTIME_DATA / "capture_log_depth_chart.pptx"


def load_points(log_path: Path) -> list[dict]:
    items = json.loads(log_path.read_text(encoding="utf-8"))

    points = []
    for idx, item in enumerate(items, start=1):
        raw = item.get("raw_depth_mm")
        depth = item.get("depth_mm")
        if raw is None or depth is None:
            continue
        raw_f = float(raw)
        depth_f = float(depth)
        points.append(
            {
                "timestamp": item.get("timestamp", ""),
                "sample_label": f"S{idx}",
                "raw_mm": raw_f,
                "depth_mm": depth_f,
                "notes": item.get("notes", ""),
            }
        )
    return points


def main() -> int:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    points = load_points(DEFAULT_LOG)
    if len(points) < 2:
        raise SystemExit("Need at least 2 rows with both raw_depth_mm and depth_mm filled.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor.from_string("F7F6F3")

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.55))
    p = title_box.text_frame.paragraphs[0]
    p.text = "거리별 Raw Depth vs 알고리즘 적용 Depth 비교"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = "NanumSquare Bold"
    p.runs[0].font.size = Pt(24)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor.from_string("17212B")

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.88), Inches(11.8), Inches(0.28))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = f"기준 로그: {DEFAULT_LOG.name} / 비교 대상: raw_depth_mm vs depth_mm"
    sp.alignment = PP_ALIGN.CENTER
    sp.runs[0].font.name = "NanumBarunGothic"
    sp.runs[0].font.size = Pt(10)
    sp.runs[0].font.color.rgb = RGBColor.from_string("5A6777")

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(1.25),
        Inches(12.3),
        Inches(5.95),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    panel.line.color.rgb = RGBColor.from_string("D8DEE8")
    panel.line.width = Pt(1.0)

    chart_data = CategoryChartData()
    chart_data.categories = [p["sample_label"] for p in points]
    chart_data.add_series("Raw depth (mm)", [p["raw_mm"] for p in points])
    chart_data.add_series("Depth mm (algorithm applied)", [p["depth_mm"] for p in points])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.85),
        Inches(1.75),
        Inches(8.45),
        Inches(4.65),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = max(max(p["raw_mm"] for p in points), max(p["depth_mm"] for p in points)) * 1.15
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "X축: 샘플 순서 / Y축: 거리(mm)"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.name = "Noto Sans CJK KR"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True

    colors = ["E85D3F", "0F9D58"]
    for series, color in zip(chart.series, colors, strict=False):
        series.format.line.color.rgb = RGBColor.from_string(color)
        series.format.line.width = Pt(2.6)
        try:
            series.marker.style = 8
            series.marker.size = 8
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = RGBColor.from_string(color)
            series.marker.format.line.color.rgb = RGBColor.from_string(color)
        except Exception:
            pass

    summary_box = slide.shapes.add_textbox(Inches(9.6), Inches(1.85), Inches(2.65), Inches(1.7))
    tf = summary_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "요약"
    p0.runs[0].font.name = "NanumSquare Bold"
    p0.runs[0].font.size = Pt(18)
    p0.runs[0].font.bold = True
    p0.runs[0].font.color.rgb = RGBColor.from_string("17212B")

    for text in [
        f"샘플 수: {len(points)}",
        "비교 항목: raw_depth_mm",
        "비교 항목: depth_mm",
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.runs[0].font.name = "NanumBarunGothic"
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = RGBColor.from_string("17212B")

    rows = len(points) + 1
    table = slide.shapes.add_table(rows, 3, Inches(9.45), Inches(3.75), Inches(3.1), Inches(2.2)).table
    table.columns[0].width = Inches(0.75)
    table.columns[1].width = Inches(1.1)
    table.columns[2].width = Inches(1.25)

    headers = ["샘플", "Raw", "Depth"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string("EAF1FB")
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.runs[0]
        run.font.name = "Noto Sans CJK KR"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("17212B")

    for row_idx, point in enumerate(points, start=1):
        values = [point["sample_label"], str(int(point["raw_mm"])), str(int(point["depth_mm"]))]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0]
            run.font.name = "NanumBarunGothic"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor.from_string("17212B")

    DEFAULT_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DEFAULT_OUT))
    print(DEFAULT_OUT)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
