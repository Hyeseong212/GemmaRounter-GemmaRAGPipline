#!/usr/bin/env python3

from __future__ import annotations

import json
from pathlib import Path


ROOT = Path("/home/rbiotech-server/LLM_Harnes_Support/GemmaRounter-GemmaRAGPipline")
RUNTIME_DATA = ROOT / "oak-wheelchair-depth-test" / "runtime-data"
DEFAULT_LOG = RUNTIME_DATA / "capture_log.json"
DEFAULT_OUT = RUNTIME_DATA / "capture_log_table.pptx"


def optional_str(value):
    if value is None:
        return "-"
    return str(value)


def make_rows(items: list[dict]) -> list[list[str]]:
    rows: list[list[str]] = []
    for idx, item in enumerate(items, start=1):
        actual = item.get("actual_distance_mm")
        depth = item.get("depth_mm")
        error = None
        if actual is not None and depth is not None:
            error = int(depth) - int(actual)

        bottom_left = item.get("pixel_coord_bottom_left") or ["-", "-"]
        rows.append(
            [
                str(idx),
                optional_str(item.get("timestamp")),
                optional_str(item.get("raw_depth_mm")),
                optional_str(depth),
                optional_str(actual),
                optional_str(error),
                f"[{bottom_left[0]}, {bottom_left[1]}]",
                optional_str(item.get("notes") or "-"),
            ]
        )
    return rows


def main() -> int:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    log_path = DEFAULT_LOG
    out_path = DEFAULT_OUT

    items = json.loads(log_path.read_text(encoding="utf-8"))
    rows = make_rows(items)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor.from_string("F7F6F3")

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Capture Log Summary Table"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = "NanumSquare Bold"
    p.runs[0].font.size = Pt(24)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor.from_string("17212B")

    sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.82), Inches(11.5), Inches(0.25))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = str(log_path)
    sp.alignment = PP_ALIGN.CENTER
    sp.runs[0].font.name = "NanumBarunGothic"
    sp.runs[0].font.size = Pt(10)
    sp.runs[0].font.color.rgb = RGBColor.from_string("5A6777")

    headers = ["No", "Timestamp", "Raw(mm)", "Depth(mm)", "Actual(mm)", "Error(mm)", "BL[x,y]", "Notes"]
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(0.35),
        Inches(1.15),
        Inches(12.65),
        Inches(5.95),
    ).table

    widths = [0.45, 1.55, 1.0, 1.0, 1.0, 0.95, 1.35, 5.35]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string("EAF1FB")
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.runs[0]
        run.font.name = "Noto Sans CJK KR"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("17212B")

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER if col_idx != 7 else PP_ALIGN.LEFT
            run = paragraph.runs[0]
            run.font.name = "NanumBarunGothic"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor.from_string("17212B")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(out_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
