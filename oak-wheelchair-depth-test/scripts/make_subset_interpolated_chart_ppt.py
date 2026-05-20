#!/usr/bin/env python3

from __future__ import annotations

import csv
import json
from pathlib import Path


ROOT = Path("/home/rbiotech-server/LLM_Harnes_Support/GemmaRounter-GemmaRAGPipline")
RUNTIME_DATA = ROOT / "oak-wheelchair-depth-test" / "runtime-data"
LOG_PATH = RUNTIME_DATA / "capture_log.json"
OUT_PPT = RUNTIME_DATA / "capture_log_depth_chart_subset_interpolated.pptx"
OUT_CSV = RUNTIME_DATA / "capture_log_depth_chart_subset_interpolated.csv"


def load_anchor_points() -> list[dict]:
    items = json.loads(LOG_PATH.read_text(encoding="utf-8"))
    valid = [
        {
            "sample": f"S{idx}",
            "raw_mm": float(item["raw_depth_mm"]),
            "depth_mm": float(item["depth_mm"]),
            "timestamp": item.get("timestamp", ""),
        }
        for idx, item in enumerate(items, start=1)
        if item.get("raw_depth_mm") is not None and item.get("depth_mm") is not None
    ]
    # User explicitly requested S4 ~ S7 only
    return [point for point in valid if point["sample"] in {"S4", "S5", "S6", "S7"}]


def interpolate(p0: dict, p1: dict, fraction: float, label: str) -> dict:
    return {
        "sample": label,
        "raw_mm": round(p0["raw_mm"] + (p1["raw_mm"] - p0["raw_mm"]) * fraction),
        "depth_mm": round(p0["depth_mm"] + (p1["depth_mm"] - p0["depth_mm"]) * fraction),
        "timestamp": "synthetic",
    }


def build_points() -> list[dict]:
    anchors = load_anchor_points()
    if len(anchors) != 4:
        raise RuntimeError("Expected exactly 4 anchor points (S4~S7) in capture_log.json")

    s4, s5, s6, s7 = anchors
    points = [
        s4,
        interpolate(s4, s5, 1 / 3, "I45-1"),
        interpolate(s4, s5, 2 / 3, "I45-2"),
        s5,
        interpolate(s5, s6, 1 / 3, "I56-1"),
        interpolate(s5, s6, 2 / 3, "I56-2"),
        s6,
        interpolate(s6, s7, 1 / 2, "I67-1"),
        s7,
    ]
    # User-requested placeholder actual-distance line:
    # keep it within about +/-100 mm of the algorithm-applied depth.
    actual_offsets = [-80, 60, -40, 70, -50, 90, -30, 50, -60]
    for point, offset in zip(points, actual_offsets, strict=True):
        point["actual_mm"] = int(point["depth_mm"] + offset)
    return points


def write_csv(points: list[dict]) -> None:
    OUT_CSV.parent.mkdir(parents=True, exist_ok=True)
    with OUT_CSV.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.DictWriter(fh, fieldnames=["sample", "raw_mm", "depth_mm", "actual_mm", "timestamp"])
        writer.writeheader()
        writer.writerows(points)


def main() -> int:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    points = build_points()
    write_csv(points)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor.from_string("F7F6F3")

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.55))
    p = title_box.text_frame.paragraphs[0]
    p.text = "S4 ~ S7 구간 Raw Depth vs Depth(mm) 비교"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = "NanumSquare Bold"
    p.runs[0].font.size = Pt(24)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor.from_string("17212B")

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.88), Inches(11.8), Inches(0.35))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "원본 anchor는 S4, S5, S6, S7이며 중간 5개 샘플은 선형 보간으로 추가"
    sp.alignment = PP_ALIGN.CENTER
    sp.runs[0].font.name = "NanumBarunGothic"
    sp.runs[0].font.size = Pt(11)
    sp.runs[0].font.color.rgb = RGBColor.from_string("5A6777")

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(1.25),
        Inches(12.3),
        Inches(5.95),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    panel.line.color.rgb = RGBColor.from_string("D8DEE8")
    panel.line.width = Pt(1.0)

    chart_data = CategoryChartData()
    chart_data.categories = [point["sample"] for point in points]
    chart_data.add_series("Raw depth (mm)", [point["raw_mm"] for point in points])
    chart_data.add_series("Depth mm", [point["depth_mm"] for point in points])
    chart_data.add_series("Actual distance (placeholder)", [point["actual_mm"] for point in points])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.85),
        Inches(1.75),
        Inches(8.55),
        Inches(4.75),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = max(
        max(point["raw_mm"] for point in points),
        max(point["depth_mm"] for point in points),
        max(point["actual_mm"] for point in points),
    ) * 1.15
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "S4~S7 구간 + 선형 보간 5개 샘플"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.name = "Noto Sans CJK KR"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True

    colors = ["E85D3F", "0F9D58", "1B6EF3"]
    for series, color in zip(chart.series, colors, strict=False):
        series.format.line.color.rgb = RGBColor.from_string(color)
        series.format.line.width = Pt(2.6)
        try:
            series.marker.style = 8
            series.marker.size = 8
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = RGBColor.from_string(color)
            series.marker.format.line.color.rgb = RGBColor.from_string(color)
        except Exception:
            pass

    summary_box = slide.shapes.add_textbox(Inches(9.7), Inches(1.85), Inches(2.55), Inches(1.95))
    tf = summary_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "구성"
    p0.runs[0].font.name = "NanumSquare Bold"
    p0.runs[0].font.size = Pt(18)
    p0.runs[0].font.bold = True
    p0.runs[0].font.color.rgb = RGBColor.from_string("17212B")
    for text in [
        "원본 샘플: 4개",
        "추가 샘플: 5개",
        "추가 방식: 선형 보간",
        "실측선: ±10cm placeholder",
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.runs[0].font.name = "NanumBarunGothic"
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = RGBColor.from_string("17212B")

    table = slide.shapes.add_table(len(points) + 1, 4, Inches(9.15), Inches(4.05), Inches(3.45), Inches(2.15)).table
    headers = ["샘플", "Raw", "Depth", "Actual"]
    widths = [0.7, 0.9, 0.95, 0.9]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string("EAF1FB")
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.runs[0]
        run.font.name = "Noto Sans CJK KR"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("17212B")
    for row_idx, point in enumerate(points, start=1):
        for col_idx, value in enumerate([point["sample"], str(int(point["raw_mm"])), str(int(point["depth_mm"])), str(int(point["actual_mm"]))]):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0]
            run.font.name = "NanumBarunGothic"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor.from_string("17212B")

    foot = slide.shapes.add_textbox(Inches(0.92), Inches(6.62), Inches(7.8), Inches(0.25))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "주: Actual distance 선은 발표용 placeholder이며, depth_mm 대비 ±10cm 이내 임의값으로 추가함."
    fp.runs[0].font.name = "NanumBarunGothic"
    fp.runs[0].font.size = Pt(10)
    fp.runs[0].font.color.rgb = RGBColor.from_string("5A6777")

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(OUT_PPT)
    print(OUT_CSV)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
