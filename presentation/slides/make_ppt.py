#!/usr/bin/env python3

from __future__ import annotations

import importlib.util
import re
from collections import Counter
from pathlib import Path

from PIL import Image


ROOT = Path("/home/rbiotech-server/LLM_Harnes_Support/GemmaRounter-GemmaRAGPipline")
PRESENTATION_ROOT = ROOT / "presentation"
EXPORT_PATH = PRESENTATION_ROOT / "export" / "gemma4_oak_obstacle_presentation_draft.pptx"
EXPORT_SHAPE_PATH = PRESENTATION_ROOT / "export" / "gemma4_oak_obstacle_presentation_editable_shapes.pptx"
FIGURES_RENDERED = PRESENTATION_ROOT / "figures" / "rendered"
ASSETS = PRESENTATION_ROOT / "assets"
USER_REPORT = Path("/home/rbiotech-server/Downloads/병원 이미지/장애물_서술형_결과.txt")

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_TITLE = "NanumSquare Bold"
FONT_BODY = "Noto Sans CJK KR"
FONT_CAPTION = "NanumBarunGothic"

BG = "F7F6F3"
PANEL = "FFFFFF"
TEXT = "17212B"
SUBTEXT = "5A6777"
ACCENT = "1B6EF3"
ACCENT_2 = "E85D3F"
ACCENT_3 = "0F9D58"
ACCENT_4 = "7C4DFF"
LINE = "D8DEE8"
SOFT_BLUE = "EAF1FB"
PLACEHOLDER_FILL = "EEF3FA"
PLACEHOLDER_LINE = "9DB4D6"

USE_IMAGE_PLACEHOLDERS = False


def run_asset_builder() -> None:
    script_path = PRESENTATION_ROOT / "figures" / "build_rendered_assets.py"
    spec = importlib.util.spec_from_file_location("build_rendered_assets", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Cannot load asset builder: {script_path}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    module.main()


def set_run_font(run, *, name: str, size: int, color: str, bold: bool = False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def set_paragraph_font(paragraph, *, name: str, size: int, color: str, bold: bool = False):
    for run in paragraph.runs:
        set_run_font(run, name=name, size=size, color=color, bold=bold)
    if not paragraph.runs:
        run = paragraph.add_run()
        run.text = paragraph.text
        paragraph.text = ""
        set_run_font(run, name=name, size=size, color=color, bold=bold)


def set_bg(slide, color: str = BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def add_textbox(slide, left: float, top: float, width: float, height: float):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = add_textbox(slide, 0.74, 0.38, 11.8, 0.72)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_font(p, name=FONT_TITLE, size=26, color=TEXT, bold=True)

    if subtitle:
        sub_box = add_textbox(slide, 1.2, 1.0, 10.9, 0.35)
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.alignment = PP_ALIGN.CENTER
        set_paragraph_font(sp, name=FONT_CAPTION, size=12, color=SUBTEXT)


def add_section_chip(slide, text: str, left: float, top: float, width: float = 1.5, color: str = ACCENT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.34),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_font(p, name=FONT_BODY, size=11, color="FFFFFF", bold=True)


def add_panel(slide, left: float, top: float, width: float, height: float, fill_color: str = PANEL):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    shape.line.color.rgb = RGBColor.from_string(LINE)
    shape.line.width = Pt(1.0)
    return shape


def add_bullets(
    slide,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 17,
    color: str = TEXT,
):
    box = add_textbox(slide, left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.space_after = Pt(8)
        p.line_spacing = 1.15
        set_paragraph_font(p, name=FONT_BODY, size=font_size, color=color)
    return box


def add_caption(slide, text: str, left: float, top: float, width: float, align=None):
    box = add_textbox(slide, left, top, width, 0.3)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT if align is None else align
    set_paragraph_font(p, name=FONT_CAPTION, size=11, color=SUBTEXT)


def add_metric_card(slide, left: float, top: float, width: float, height: float, label: str, value: str, value_color: str):
    add_panel(slide, left, top, width, height, fill_color=SOFT_BLUE)
    label_box = add_textbox(slide, left + 0.08, top + 0.1, width - 0.16, 0.22)
    lp = label_box.text_frame.paragraphs[0]
    lp.text = label
    lp.alignment = PP_ALIGN.CENTER
    set_paragraph_font(lp, name=FONT_CAPTION, size=11, color=SUBTEXT, bold=True)

    value_box = add_textbox(slide, left + 0.08, top + 0.35, width - 0.16, height - 0.4)
    vp = value_box.text_frame.paragraphs[0]
    vp.text = value
    vp.alignment = PP_ALIGN.CENTER
    set_paragraph_font(vp, name=FONT_TITLE, size=18, color=value_color, bold=True)


def add_line_rect(slide, left: float, top: float, width: float, height: float, color: str, rotation: float = 0.0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    shape.rotation = rotation
    return shape


def add_dot(slide, center_x: float, center_y: float, diameter: float, color: str):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(center_x - diameter / 2),
        Inches(center_y - diameter / 2),
        Inches(diameter),
        Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    return shape


def add_label_pill(slide, text: str, left: float, top: float, width: float, height: float = 0.38, font_size: int = 12):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    shape.line.color.rgb = RGBColor.from_string(LINE)
    shape.line.width = Pt(1.0)
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_font(p, name=FONT_CAPTION, size=font_size, color=TEXT)
    return shape


def add_down_arrow(slide, center_x: float, top: float, width: float, height: float, color: str):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
        Inches(center_x - width / 2),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    return shape


def parse_report_counts() -> list[tuple[str, int]]:
    if not USER_REPORT.exists():
        return []
    text = USER_REPORT.read_text(encoding="utf-8")
    classes = re.findall(r"클래스: ([a-zA-Z_]+)", text)
    return Counter(classes).most_common(6)


def add_image_placeholder(
    slide,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    title: str | None = None,
):
    add_panel(slide, left, top, width, height, fill_color=PLACEHOLDER_FILL)
    inner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.18),
        Inches(top + 0.18),
        Inches(width - 0.36),
        Inches(height - 0.36),
    )
    inner.fill.solid()
    inner.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    inner.line.color.rgb = RGBColor.from_string(PLACEHOLDER_LINE)
    inner.line.width = Pt(1.6)

    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.28),
        Inches(top + 0.28),
        Inches(min(width - 0.56, 1.6)),
        Inches(0.34),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor.from_string(ACCENT)
    chip.line.fill.background()
    cp = chip.text_frame.paragraphs[0]
    cp.text = "EDITABLE SHAPE"
    cp.alignment = PP_ALIGN.CENTER
    set_paragraph_font(cp, name=FONT_BODY, size=10, color="FFFFFF", bold=True)

    name = image_path.name
    dims = ""
    if image_path.exists():
        try:
            with Image.open(image_path) as img:
                dims = f"{img.size[0]} x {img.size[1]}"
        except Exception:
            dims = ""

    title_box = add_textbox(slide, left + 0.34, top + 0.78, width - 0.68, 0.46)
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title or name
    tp.alignment = PP_ALIGN.CENTER
    set_paragraph_font(tp, name=FONT_TITLE, size=18, color=TEXT, bold=True)

    body_box = add_textbox(slide, left + 0.45, top + 1.45, width - 0.9, height - 1.8)
    tf = body_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"파일명: {name}"
    p1.alignment = PP_ALIGN.CENTER
    set_paragraph_font(p1, name=FONT_BODY, size=14, color=TEXT)
    if dims:
        p2 = tf.add_paragraph()
        p2.text = f"원본 크기: {dims}"
        p2.alignment = PP_ALIGN.CENTER
        set_paragraph_font(p2, name=FONT_BODY, size=13, color=SUBTEXT)
    p3 = tf.add_paragraph()
    p3.text = "이 영역은 직접 편집 가능한 도형 placeholder입니다."
    p3.alignment = PP_ALIGN.CENTER
    set_paragraph_font(p3, name=FONT_CAPTION, size=12, color=SUBTEXT)
    icon_box = add_textbox(slide, left + width / 2 - 0.6, top + height - 1.1, 1.2, 0.35)
    ip = icon_box.text_frame.paragraphs[0]
    ip.text = "[ IMAGE ]"
    ip.alignment = PP_ALIGN.CENTER
    set_paragraph_font(ip, name=FONT_BODY, size=11, color=PLACEHOLDER_LINE, bold=True)


def add_image_contain(slide, image_path: Path, left: float, top: float, width: float, height: float, pad: float = 0.0):
    if USE_IMAGE_PLACEHOLDERS:
        return add_image_placeholder(slide, image_path, left, top, width, height)
    if not image_path.exists():
        return None
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    box_w = max(width - pad * 2, 0.1)
    box_h = max(height - pad * 2, 0.1)
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        draw_w = box_w
        draw_h = box_w / img_ratio
    else:
        draw_h = box_h
        draw_w = box_h * img_ratio
    x = left + (width - draw_w) / 2
    y = top + (height - draw_h) / 2
    return slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(draw_w), height=Inches(draw_h))


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor.from_string(ACCENT)
    band.line.fill.background()

    add_title(slide, "OAK 카메라와 Gemma 4를 이용한\n병원 장애물 인식 및 위치 추론", "환자이송 로봇 적용을 위한 이미지 추론·좌표 후처리·거리 측정")

    add_panel(slide, 0.8, 1.65, 4.7, 1.35)
    add_bullets(
        slide,
        [
            "Gemma 4 기반 병원 장애물 클래스 분류",
            "좌표 후처리와 OAK depth 결합",
        ],
        1.0,
        1.92,
        4.2,
        0.9,
        font_size=16,
    )

    add_panel(slide, 0.8, 3.25, 4.7, 2.45)
    add_bullets(
        slide,
        [
            "병실, 병동, 복도, 로비 환경을 이미지로 입력했다.",
            "환자이송 로봇 주행에 필요한 장애물 위치와 거리를 추론했다.",
            "실제 휠체어 거리 측정 사례까지 확인했다.",
        ],
        1.0,
        3.5,
        4.2,
        1.9,
        font_size=15,
    )

    add_panel(slide, 5.8, 1.55, 6.75, 4.85)
    add_image_contain(slide, FIGURES_RENDERED / "05_user_samples_contact_sheet.png", 6.0, 1.75, 6.35, 4.25)
    add_caption(slide, "직접 테스트한 병원 샘플 이미지 묶음", 6.0, 6.1, 4.0)


def slide_challenge_frame(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "문제 구조화: 병원 장애물 인식이 왜 어려운가", "AS-IS / TO-BE 관점에서 병원 환경의 복잡도를 정리")
    add_section_chip(slide, "Challenge", 0.72, 0.18, width=1.55, color=ACCENT_2)
    add_panel(slide, 0.48, 1.28, 12.35, 5.92)

    add_line_rect(slide, 1.05, 6.82, 11.0, 0.04, TEXT)
    add_line_rect(slide, 1.05, 1.95, 0.04, 4.91, TEXT)
    add_caption(slide, "환경 복잡도", 11.1, 6.92, 1.2)
    axis_y = add_textbox(slide, 0.18, 2.0, 0.7, 0.35)
    yp = axis_y.text_frame.paragraphs[0]
    yp.text = "자율성"
    set_paragraph_font(yp, name=FONT_BODY, size=13, color=TEXT, bold=True)

    add_line_rect(slide, 3.95, 2.25, 0.015, 4.57, LINE)
    add_line_rect(slide, 7.2, 2.25, 0.015, 4.57, LINE)

    add_panel(slide, 1.42, 5.12, 2.4, 1.42, fill_color="EEF2F7")
    asis_title = add_textbox(slide, 2.02, 5.3, 1.2, 0.28)
    ap = asis_title.text_frame.paragraphs[0]
    ap.text = "AS-IS"
    ap.alignment = PP_ALIGN.CENTER
    set_paragraph_font(ap, name=FONT_TITLE, size=20, color=TEXT, bold=True)
    asis_body = add_textbox(slide, 1.72, 5.66, 1.8, 0.58)
    asis_tf = asis_body.text_frame
    asis_tf.word_wrap = True
    p1 = asis_tf.paragraphs[0]
    p1.text = "장애물 존재 여부 중심"
    p1.alignment = PP_ALIGN.CENTER
    set_paragraph_font(p1, name=FONT_BODY, size=15, color=TEXT)
    p2 = asis_tf.add_paragraph()
    p2.text = "규칙 기반 회피"
    p2.alignment = PP_ALIGN.CENTER
    set_paragraph_font(p2, name=FONT_BODY, size=15, color=TEXT)

    add_panel(slide, 8.82, 2.02, 3.5, 2.18, fill_color="EAF1FB")
    tobe_title = add_textbox(slide, 9.95, 2.22, 1.25, 0.3)
    tp = tobe_title.text_frame.paragraphs[0]
    tp.text = "TO-BE"
    tp.alignment = PP_ALIGN.CENTER
    set_paragraph_font(tp, name=FONT_TITLE, size=20, color=ACCENT, bold=True)
    tobe_body = add_textbox(slide, 9.18, 2.72, 2.85, 1.15)
    tobe_tf = tobe_body.text_frame
    tobe_tf.word_wrap = True
    for idx, text in enumerate(["병원 장애물 의미 이해", "좌표 / 거리 기반 행동 판단", "양보 · 대기 · 우회 의사결정"]):
        p = tobe_tf.paragraphs[0] if idx == 0 else tobe_tf.add_paragraph()
        p.text = f"•  {text}"
        set_paragraph_font(p, name=FONT_BODY, size=15, color=TEXT)
        p.space_after = Pt(4)

    # trajectory
    add_line_rect(slide, 4.0, 4.55, 5.15, 0.05, "C7771A", rotation=-24)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(8.9), Inches(3.15), Inches(0.28), Inches(0.24))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor.from_string("C7771A")
    arrow.line.fill.background()
    arrow.rotation = -24

    challenge_points = [
        (1.22, 5.85, ACCENT, "정적 복도 / 단순 주행", 1.02, 6.33, 1.82),
        (3.45, 5.48, ACCENT_2, "병실 진입 / 좁은 회전", 2.72, 5.94, 1.95),
        (5.42, 4.92, ACCENT_3, "복도 혼잡 / 사람 · 카트 공존", 4.32, 5.48, 2.42),
        (7.32, 4.18, "8E44AD", "장애물 분류 + 위치 추론", 6.7, 4.72, 2.0),
        (8.95, 3.02, "F39C12", "양보 · 대기 · 우회 판단", 8.0, 4.05, 1.92),
    ]
    for cx, cy, color, label, lx, ly, lw in challenge_points:
        add_dot(slide, cx, cy, 0.28, color)
        add_label_pill(slide, label, lx, ly, lw, height=0.42, font_size=12)

    add_panel(slide, 8.35, 5.02, 3.95, 1.72)
    box_head = add_textbox(slide, 8.72, 5.28, 1.65, 0.28)
    bh = box_head.text_frame.paragraphs[0]
    bh.text = "핵심 도전 요소"
    set_paragraph_font(bh, name=FONT_BODY, size=16, color=ACCENT_2, bold=True)
    add_bullets(
        slide,
        [
            "복도 / 병실 환경은 비정형적이다.",
            "사람과 의료 장비가 계속 바뀐다.",
            "단순 회피가 아니라 의미 기반 대응이 필요하다.",
        ],
        8.72,
        5.64,
        3.18,
        0.98,
        font_size=15,
    )


def slide_framework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "전체 목표 프레임워크", "인지 -> 의미 이해 -> 공간 정렬 -> 행동 판단으로 연결되는 구조")
    add_section_chip(slide, "Framework", 0.72, 0.18, width=1.6, color=ACCENT_3)
    add_panel(slide, 0.6, 1.35, 12.1, 5.7)

    layers = [
        (1.0, 1.95, 11.0, 0.82, "1. 인지 계층", "OAK RGB / Depth, 병원 장면 입력, 프레임 수집", ACCENT, "EAF1FB"),
        (1.45, 3.02, 10.1, 0.82, "2. 의미 이해 계층", "Gemma 4 기반 장애물 클래스 추론, 장면 의미 해석", ACCENT_2, "FFF4EE"),
        (1.9, 4.09, 9.2, 0.82, "3. 공간 정렬 계층", "normalized [y,x] -> 픽셀 -> bottom-left 수학좌표 + ROI depth", ACCENT_3, "EEF7F0"),
        (2.35, 5.16, 8.3, 0.82, "4. 행동 판단 계층", "양보 / 대기 / 우회 / 음성 안내 / 로봇 주행 입력", ACCENT_4, "F3EEFB"),
    ]
    for left, top, width, height, title, subtitle, color, fill in layers:
        shape = add_panel(slide, left, top, width, height, fill_color=fill)
        shape.line.color.rgb = RGBColor.from_string(color)
        shape.line.width = Pt(1.8)
        tbox = add_textbox(slide, left + 0.2, top + 0.14, width - 0.4, 0.22)
        tp = tbox.text_frame.paragraphs[0]
        tp.text = title
        set_paragraph_font(tp, name=FONT_TITLE, size=17, color=color, bold=True)
        sbox = add_textbox(slide, left + 0.2, top + 0.43, width - 0.4, 0.2)
        sp = sbox.text_frame.paragraphs[0]
        sp.text = subtitle
        set_paragraph_font(sp, name=FONT_BODY, size=13, color=TEXT)

    add_down_arrow(slide, 6.5, 2.76, 0.22, 0.28, ACCENT)
    add_down_arrow(slide, 6.5, 3.83, 0.22, 0.28, ACCENT_2)
    add_down_arrow(slide, 6.5, 4.9, 0.22, 0.28, ACCENT_3)

    add_panel(slide, 10.15, 5.12, 2.1, 0.92)
    t = add_textbox(slide, 10.35, 5.28, 1.7, 0.18)
    p = t.text_frame.paragraphs[0]
    p.text = "예시 로봇 멘트"
    set_paragraph_font(p, name=FONT_BODY, size=13, color=ACCENT_4, bold=True)
    c = add_textbox(slide, 10.35, 5.54, 1.6, 0.28)
    cp = c.text_frame.paragraphs[0]
    cp.text = "“전방에 의료진이 있습니다.\n잠시 대기 후 이동하겠습니다.”"
    set_paragraph_font(cp, name=FONT_CAPTION, size=11, color=TEXT)


def slide_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "사용한 기술 설명", "하드웨어, 서버 추론, 후처리, 거리 측정")
    add_section_chip(slide, "Tech Stack", 0.72, 0.18, width=1.6)

    card_specs = [
        (0.8, 1.55, 2.95, 4.8, ACCENT, "OAK Camera", ["RGB + Depth 동시 수집", "Stereo depth 기반 거리 조회", "실측 휠체어 depth 확인"]),
        (3.98, 1.55, 2.95, 4.8, ACCENT_2, "Gemma 4", ["이미지 기반 장애물 클래스 추론", "위치 좌표를 정규화 값으로 출력", "프롬프트로 taxonomy 조정 가능"]),
        (7.16, 1.55, 2.95, 4.8, ACCENT_3, "Post-processing", ["[y,x] -> 픽셀 좌표 변환", "top-left -> bottom-left 수학좌표 변환", "좌표를 로봇 사용 형식으로 정리"]),
        (10.34, 1.55, 2.2, 4.8, ACCENT_4, "Runtime", ["C++ 테스트 하네스", "18088 /infer 서버 연동", "배치 테스트 및 PPT 자동 생성"]),
    ]
    for left, top, width, height, color, title, bullets in card_specs:
        add_panel(slide, left, top, width, height)
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.18), Inches(top + 0.18), Inches(min(width - 0.36, 1.45)), Inches(0.38))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor.from_string(color)
        chip.line.fill.background()
        cp = chip.text_frame.paragraphs[0]
        cp.text = title
        cp.alignment = PP_ALIGN.CENTER
        set_paragraph_font(cp, name=FONT_BODY, size=11, color="FFFFFF", bold=True)
        add_bullets(slide, bullets, left + 0.18, top + 0.8, width - 0.36, height - 1.0, font_size=14)


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "파이프라인 다이어그램", "입력 이미지에서 클래스·좌표·거리 정보까지 생성")
    add_section_chip(slide, "Pipeline", 0.72, 0.18)
    add_panel(slide, 0.58, 1.35, 12.15, 5.55)
    steps = [
        (0.95, 2.45, 2.2, 1.1, "1. OAK Camera", "RGB / Depth 입력", ACCENT),
        (3.35, 2.45, 2.5, 1.1, "2. Gemma 4 /infer", "장애물 클래스 + 위치 추론", ACCENT_2),
        (6.1, 2.45, 2.6, 1.1, "3. 좌표 후처리", "normalized -> pixel -> math", ACCENT_3),
        (8.95, 2.45, 2.6, 1.1, "4. Depth ROI", "median depth 계산", ACCENT_4),
    ]
    for left, top, width, height, title, subtitle, color in steps:
        shape = add_panel(slide, left, top, width, height, fill_color="FFFFFF")
        shape.line.color.rgb = RGBColor.from_string(color)
        shape.line.width = Pt(1.6)
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.12), Inches(top + 0.12), Inches(1.1), Inches(0.28))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor.from_string(color)
        chip.line.fill.background()
        p = chip.text_frame.paragraphs[0]
        p.text = title.split(". ", 1)[0]
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_font(p, name=FONT_BODY, size=10, color="FFFFFF", bold=True)
        t = add_textbox(slide, left + 0.16, top + 0.46, width - 0.3, 0.22)
        tp = t.text_frame.paragraphs[0]
        tp.text = title
        set_paragraph_font(tp, name=FONT_TITLE, size=16, color=TEXT, bold=True)
        s = add_textbox(slide, left + 0.16, top + 0.75, width - 0.3, 0.18)
        sp = s.text_frame.paragraphs[0]
        sp.text = subtitle
        set_paragraph_font(sp, name=FONT_BODY, size=12, color=SUBTEXT)

    for x in [3.0, 5.85, 8.7]:
        chevron = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.78), Inches(0.28), Inches(0.42))
        chevron.fill.solid()
        chevron.fill.fore_color.rgb = RGBColor.from_string(ACCENT)
        chevron.line.fill.background()

    add_panel(slide, 1.0, 4.25, 10.7, 1.35)
    add_bullets(
        slide,
        [
            "OAK에서 받은 병원 장면을 Gemma 4로 보내 장애물 클래스와 위치를 추론한다.",
            "위치 정보는 후처리로 로봇 주행에 쓸 수 있는 bottom-left 수학좌표로 변환한다.",
            "필요 시 동일 지점 depth ROI median을 계산해 거리까지 결합한다.",
        ],
        1.25,
        4.52,
        10.1,
        0.8,
        font_size=15,
    )


def slide_actual_inference(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "실제 동작 1: 이미지 추론 결과", "직접 테스트한 병원 샘플과 추론 결과 분포")
    add_section_chip(slide, "Actual Run", 0.72, 0.18, width=1.55)

    add_panel(slide, 0.7, 1.45, 5.95, 4.2)
    add_image_contain(slide, FIGURES_RENDERED / "05_user_samples_contact_sheet.png", 0.88, 1.66, 5.58, 3.55)
    add_caption(slide, "직접 수집한 병원 이미지 샘플", 0.92, 5.38, 2.8)

    add_panel(slide, 6.82, 1.45, 5.82, 4.2)
    counts = parse_report_counts()
    if counts:
        max_v = max(v for _, v in counts)
        chart_left = 7.12
        base_y = 4.72
        bar_h_max = 1.82
        bar_w = 0.48
        gap = 0.18
        palette = [ACCENT, ACCENT_2, ACCENT_3, ACCENT_4, "F39C12", "00897B"]
        add_line_rect(slide, chart_left, base_y, 4.82, 0.015, TEXT)
        add_line_rect(slide, chart_left, 2.02, 0.015, 2.71, TEXT)
        for idx, (label, value) in enumerate(counts[:6]):
            bar_left = chart_left + 0.38 + idx * (bar_w + gap)
            bar_h = (value / max_v) * bar_h_max
            add_panel(slide, bar_left, base_y - bar_h, bar_w, bar_h, fill_color=palette[idx % len(palette)])
            vbox = add_textbox(slide, bar_left - 0.02, base_y - bar_h - 0.28, bar_w + 0.05, 0.18)
            vp = vbox.text_frame.paragraphs[0]
            vp.text = str(value)
            vp.alignment = PP_ALIGN.CENTER
            set_paragraph_font(vp, name=FONT_BODY, size=11, color=TEXT, bold=True)
            lbox = add_textbox(slide, bar_left - 0.2, 4.82, bar_w + 0.45, 0.52)
            lp = lbox.text_frame.paragraphs[0]
            lp.text = label
            lp.alignment = PP_ALIGN.CENTER
            set_paragraph_font(lp, name=FONT_CAPTION, size=10, color=TEXT)
        add_caption(slide, "장애물 클래스 빈도 분포", 7.02, 5.38, 2.6)

        rows = min(4, len(counts)) + 1
        table = slide.shapes.add_table(rows, 2, Inches(9.95), Inches(1.78), Inches(2.25), Inches(2.18)).table
        table.columns[0].width = Inches(1.5)
        table.columns[1].width = Inches(0.75)
        headers = [("클래스", 0), ("빈도", 1)]
        for text, col in headers:
            cell = table.cell(0, col)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(SOFT_BLUE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            set_paragraph_font(p, name=FONT_BODY, size=11, color=TEXT, bold=True)
        for row_idx, (label, value) in enumerate(counts[:4], start=1):
            c0 = table.cell(row_idx, 0)
            c1 = table.cell(row_idx, 1)
            c0.text = label
            c1.text = str(value)
            for cell in (c0, c1):
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                set_paragraph_font(p, name=FONT_CAPTION, size=10, color=TEXT)
        add_caption(slide, "상위 클래스 요약", 10.0, 4.15, 1.8)

    add_panel(slide, 0.7, 5.95, 11.94, 0.72)
    add_bullets(
        slide,
        [
            "사용자 샘플 이미지 8장을 입력해 장애물 클래스와 위치를 서술형 보고서로 정리했다.",
            "병실, 복도, 병동, 로비에서 wheelchair, medical_cart, iv_pole, person 등 주요 객체가 반복적으로 검출됐다.",
        ],
        0.92,
        6.13,
        11.4,
        0.42,
        font_size=14,
    )


def slide_actual_coordinates(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "실제 동작 2: 좌표 후처리와 병실 데모", "모델 좌표를 로봇 사용 좌표로 변환")
    add_section_chip(slide, "Coordinate", 0.72, 0.18, width=1.6, color=ACCENT_2)

    add_panel(slide, 0.7, 1.45, 4.95, 4.8)
    add_image_contain(slide, FIGURES_RENDERED / "02_coordinate_transform.png", 0.9, 1.65, 4.55, 4.1)
    add_caption(slide, "정규화 좌표 -> 픽셀 -> bottom-left 수학좌표", 0.92, 5.95, 3.7)

    add_panel(slide, 5.9, 1.45, 6.75, 4.8)
    add_image_contain(slide, FIGURES_RENDERED / "04_shared_room_overlay.png", 6.08, 1.62, 6.38, 4.15)
    add_caption(slide, "병실 샘플에서 주요 장애물 좌표를 시각화", 6.12, 5.95, 3.7)

    add_metric_card(slide, 0.9, 6.15, 1.85, 0.72, "대표 클래스", "wheelchair", ACCENT)
    add_metric_card(slide, 2.95, 6.15, 1.7, 0.72, "좌표 형식", "[x, y]", ACCENT_2)
    add_metric_card(slide, 5.95, 6.15, 1.65, 0.72, "Detected", "5", ACCENT_3)
    add_metric_card(slide, 7.8, 6.15, 2.65, 0.72, "활용 출력", "주행용 위치 입력", ACCENT_4)


def slide_actual_depth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "실제 동작 3: 휠체어 거리 측정", "좌표와 OAK depth를 결합한 실측 사례")
    add_section_chip(slide, "Depth", 0.72, 0.18)

    add_panel(slide, 0.7, 1.45, 7.95, 5.15)
    add_image_contain(slide, ASSETS / "장애물거리측정.png", 0.88, 1.62, 7.58, 4.65)
    add_caption(slide, "실제 병원 복도 이미지에서 휠체어 좌표와 depth를 동시에 표기", 0.92, 6.28, 4.5)

    add_panel(slide, 8.9, 1.45, 3.75, 5.15)
    add_bullets(
        slide,
        [
            "Gemma 4가 휠체어 위치를 추론한다.",
            "후처리에서 bottom-left와 top-left 좌표를 모두 계산한다.",
            "동일 지점 ROI depth median으로 거리 값을 얻는다.",
        ],
        9.12,
        1.82,
        3.3,
        1.85,
        font_size=15,
    )
    add_metric_card(slide, 9.18, 4.1, 3.1, 0.72, "Bottom-left", "BL [451, 199]", ACCENT)
    add_metric_card(slide, 9.18, 4.92, 3.1, 0.72, "Top-left", "TL [451, 521]", ACCENT_2)
    add_metric_card(slide, 9.18, 5.74, 3.1, 0.72, "Measured depth", "2827 mm", ACCENT_3)


def slide_use_cases(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "이 기술이 어디에서 쓰이면 좋은가", "환자이송 로봇의 실제 운영 시나리오")
    add_section_chip(slide, "Use Cases", 0.72, 0.18, width=1.55, color=ACCENT_3)

    scenarios = [
        (
            0.78,
            ACCENT,
            "병실 출입 및 퇴실",
            [
                "전방 휠체어와 IV pole을 인식하고 안전거리 확보 후 진입한다.",
                "병실 내부 침대 주변 가구를 보고 회전 반경을 보수적으로 잡는다.",
            ],
            "예시 안내: “전방에 휠체어가 있습니다. 잠시 경로를 조정하겠습니다.”",
        ),
        (
            4.48,
            ACCENT_2,
            "복도 주행 및 양보",
            [
                "복도에서 의료진, 보호자, 이동식 카트를 구분해 우선순위를 판단한다.",
                "사람이 지나가는 동안 일시 정지하거나 벽 쪽으로 붙는 동작에 활용할 수 있다.",
            ],
            "예시 안내: “파란 가운을 입은 의료진이 지나가는 중입니다. 잠시 대기 후 이동하겠습니다.”",
        ),
        (
            8.18,
            ACCENT_3,
            "중앙 관제·재경로",
            [
                "장애물 클래스와 거리 정보를 서버로 보내 동선 혼잡도를 기록할 수 있다.",
                "병동별 잦은 장애물 패턴을 보고 재경로 정책을 보강할 수 있다.",
            ],
            "예시 안내: “전방 2.8m 지점에 이동 장애물이 있습니다. 우회 경로를 탐색합니다.”",
        ),
    ]

    for left, color, title, bullets, speech in scenarios:
        add_panel(slide, left, 1.55, 3.45, 4.95)
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.18), Inches(1.75), Inches(1.45), Inches(0.36))
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor.from_string(color)
        chip.line.fill.background()
        cp = chip.text_frame.paragraphs[0]
        cp.text = title
        cp.alignment = PP_ALIGN.CENTER
        set_paragraph_font(cp, name=FONT_BODY, size=11, color="FFFFFF", bold=True)
        add_bullets(slide, bullets, left + 0.18, 2.28, 3.05, 1.95, font_size=14)
        speech_box = add_textbox(slide, left + 0.18, 4.62, 3.05, 1.48)
        tf = speech_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = speech
        p.alignment = PP_ALIGN.LEFT
        set_paragraph_font(p, name=FONT_CAPTION, size=12, color=SUBTEXT)


def slide_limitations_next(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "한계와 다음 단계", "현재 기준선에서 추가로 보강할 부분")
    add_section_chip(slide, "Next Step", 0.72, 0.18, width=1.6)

    add_panel(slide, 0.8, 1.55, 5.8, 4.95)
    add_panel(slide, 6.75, 1.55, 5.8, 4.95)

    left_head = add_textbox(slide, 1.05, 1.85, 2.8, 0.35)
    lp = left_head.text_frame.paragraphs[0]
    lp.text = "현재 한계"
    set_paragraph_font(lp, name=FONT_TITLE, size=20, color=TEXT, bold=True)

    add_bullets(
        slide,
        [
            "VLM 기반이라 누락과 과추론 가능성이 있다.",
            "장면이 복잡하면 객체 수를 완전히 세지 못할 수 있다.",
            "정밀 bbox detector 수준의 좌표는 아니다.",
        ],
        1.05,
        2.38,
        5.1,
        2.15,
        font_size=16,
    )

    right_head = add_textbox(slide, 7.0, 1.85, 2.8, 0.35)
    rp = right_head.text_frame.paragraphs[0]
    rp.text = "다음 단계"
    set_paragraph_font(rp, name=FONT_TITLE, size=20, color=TEXT, bold=True)

    add_bullets(
        slide,
        [
            "detector/grounding 모델을 병행해 좌표 정밀도를 높인다.",
            "장애물 taxonomy를 환자이송 로봇 운영 기준으로 더 정제한다.",
            "OAK depth를 실시간 주행 제어와 연결한다.",
        ],
        7.0,
        2.38,
        5.1,
        2.15,
        font_size=16,
    )


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor.from_string(ACCENT)
    band.line.fill.background()

    box = add_textbox(slide, 1.05, 1.4, 11.2, 3.8)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Gemma 4와 OAK를 이용해\n병원 장애물을 인식하고 위치·거리까지 연결했습니다."
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_font(p, name=FONT_TITLE, size=28, color="FFFFFF", bold=True)

    sp = tf.add_paragraph()
    sp.text = "발표 시점 기준으로 이미지 기반 장애물 분류, 좌표 후처리, 실제 depth 측정까지 하나의 흐름으로 검증했습니다."
    sp.alignment = PP_ALIGN.CENTER
    sp.space_before = Pt(22)
    set_paragraph_font(sp, name=FONT_BODY, size=18, color="FFFFFF")


def build_presentation(prs):
    slide_title(prs)
    slide_challenge_frame(prs)
    slide_framework(prs)
    slide_tech_stack(prs)
    slide_pipeline(prs)
    slide_actual_inference(prs)
    slide_actual_coordinates(prs)
    slide_actual_depth(prs)
    slide_use_cases(prs)
    slide_limitations_next(prs)
    slide_closing(prs)


def create_presentation(export_path: Path, *, use_image_placeholders: bool):
    global USE_IMAGE_PLACEHOLDERS
    USE_IMAGE_PLACEHOLDERS = use_image_placeholders
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_presentation(prs)
    export_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(export_path))
    print(export_path)


def main() -> int:
    run_asset_builder()

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit("python-pptx가 필요합니다. 먼저 `python -m pip install python-pptx` 를 실행하세요.") from exc

    globals().update(
        RGBColor=RGBColor,
        MSO_AUTO_SHAPE_TYPE=MSO_AUTO_SHAPE_TYPE,
        PP_ALIGN=PP_ALIGN,
        Inches=Inches,
        Pt=Pt,
        Presentation=Presentation,
    )
    create_presentation(EXPORT_PATH, use_image_placeholders=False)
    create_presentation(EXPORT_SHAPE_PATH, use_image_placeholders=True)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
